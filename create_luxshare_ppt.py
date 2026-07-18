from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x4F, 0xC3, 0xF7)
GREEN = RGBColor(0x4E, 0xCC, 0xA3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x0F, 0x0F, 0x1A)
GRAY = RGBColor(0xA0, 0xA0, 0xB8)
DARK_CARD = RGBColor(0x1A, 0x1A, 0x2E)
SUBDUED = RGBColor(0x7A, 0x7A, 0x9A)
ORANGE = RGBColor(0xFF, 0xA7, 0x26)
RED = RGBColor(0xEF, 0x53, 0x50)

def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK
def add_shape(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s
def add_text(slide, l, t, w, h, tx, sz=14, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tx; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.alignment = a
    return tb
def add_multi(slide, l, t, w, h, lines, sz=14, c=GRAY, bf=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz); p.font.color.rgb = c
        if bf and i==0: p.font.bold = True
        p.space_after = Pt(4)
    return tb

# SLIDE 1: Cover
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
add_shape(s, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), BLUE)
add_text(s, Inches(1), Inches(1.5), Inches(11), Inches(1), "LUXSHARE ICT", 60, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.4), Inches(11), Inches(0.6), "Financial Analysis  FY2025", 32, BLUE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.4), Inches(11), Inches(0.5), "SHE: 002475  |  Luxshare Precision Industry Co., Ltd.", 16, GRAY, False, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.4), "Global Precision Manufacturing  |  Apple Supply Chain  |  Auto Tier-1  |  Fortune Global 500", 13, SUBDUED, False, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.4), "Revenue RMB 332.3B  |  Market Cap ~$239B  |  HKD IPO Pending", 14, GREEN, False, PP_ALIGN.CENTER)

# SLIDE 2: Executive Summary
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), BLUE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Executive Summary", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A,0x2A,0x4A))
add_multi(s, Inches(0.5), Inches(2.0), Inches(12), Inches(4.5), [
    "Luxshare ICT delivered record revenue of RMB 332.3B (+23.6% YoY) and net profit of RMB 16.6B (+24.2% YoY) in FY2025, driven by explosive automotive growth from the Leoni acquisition and AI infrastructure demand.",
    "",
    "KEY HIGHLIGHTS:",
    "  Revenue: RMB 332.3B (+23.6%)        |   Net Profit: RMB 16.6B (+24.2%)",
    "  Automotive: RMB 39.3B (+185%)        |   Data Center: RMB 24.6B (+33.8%)",
    "  Consumer Electronics: ~RMB 264.5B (+13.4%)  |   ROE: 19.1%",
    "  Cash: RMB 60.6B                     |   P/E: 17.1x (Forward ~14.5x)",
    "",
    "CEO Wang Laisheng: \"The true strength of an enterprise is demonstrated not only by its ability to grow in favorable cycles, but also by its capacity to sustain performance amid uncertainty.\"",
], 13, GRAY)

# SLIDE 3: Income Statement
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), BLUE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Income Statement", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A,0x2A,0x4A))
tbl_data = [
    ("Metric (RMB M)", "FY2025", "FY2024", "Change", "Trend"),
    ("Operating Revenue", "332,344", "268,795", "+23.6%", "Up"),
    ("Gross Profit", "38,444", "27,070", "+42.0%", "Up"),
    ("Gross Margin", "11.57%", "10.07%", "+150 bps", "Up"),
    ("EBITDA", "26,255", "21,104", "+24.4%", "Up"),
    ("EBIT", "13,653", "10,513", "+29.9%", "Up"),
    ("Net Profit (Attr.)", "16,600", "13,366", "+24.2%", "Up"),
    ("Core Net Profit", "14,170", "11,700", "+21.2%", "Up"),
    ("EPS (Diluted)", "RMB 2.26", "RMB 1.85", "+22.2%", "Up"),
]
tbl = s.shapes.add_table(len(tbl_data),5,Inches(0.5),Inches(2.0),Inches(12.3),Inches(4.2)).table
for ci in range(5): tbl.columns[ci].width = Inches(12.3/5)
for ri,row in enumerate(tbl_data):
    for ci,val in enumerate(row):
        c = tbl.cell(ri,ci); c.text = val
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(10.5); p.alignment = PP_ALIGN.CENTER
            if ri==0: p.font.bold=True; p.font.color.rgb=BLUE
            else: p.font.color.rgb=WHITE
        c.fill.solid(); c.fill.fore_color.rgb = DARK_CARD if ri%2==0 else RGBColor(0x16,0x16,0x2A)

# SLIDE 4: KPI Cards
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), BLUE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Key Performance Indicators", 28, WHITE, True)
kpis = [
    ("Revenue FY2025", "RMB 332.3B", "+23.6% YoY", GREEN),
    ("Net Profit", "RMB 16.6B", "+24.2% YoY", GREEN),
    ("Automotive", "RMB 39.3B", "+185% (Leoni)", GREEN),
    ("Data Center", "RMB 24.6B", "+33.8% (AI)", BLUE),
    ("Gross Margin", "11.57%", "+150 bps YoY", GREEN),
    ("ROE", "19.1%", "Strong returns", GREEN),
    ("Cash", "RMB 60.6B", "Healthy liquidity", BLUE),
    ("P/E (Fwd)", "~14.5x", "GS target RMB 106", ORANGE),
]
for i,(lbl,val,chg,clr) in enumerate(kpis):
    col=i%4; row=i//4
    x=Inches(0.5+col*3.15); y=Inches(2.0+row*2.2)
    card=add_shape(s,x,y,Inches(2.9),Inches(1.3),DARK_CARD)
    add_text(s,x+Inches(0.2),y+Inches(0.1),Inches(2.5),Inches(0.3),lbl,10,SUBDUED)
    add_text(s,x+Inches(0.2),y+Inches(0.35),Inches(2.5),Inches(0.5),val,22,WHITE,True)
    add_text(s,x+Inches(0.2),y+Inches(0.85),Inches(2.5),Inches(0.3),chg,11,clr)

# SLIDE 5: Segments
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), BLUE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Segment Performance", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A,0x2A,0x4A))
seg_data = [
    ("Segment", "Revenue (RMB B)", "YoY Growth", "Gross Margin", "Key Driver"),
    ("Consumer Electronics", "~264.5", "+13.4%", "~10%", "AI devices, Apple, ODM"),
    ("Automotive Electronics", "39.25", "+185.3%", "~15%*", "Leoni acquisition, Tier-1"),
    ("Comms & Data Center", "24.57", "+33.8%", "18.4%", "AI infra, high-speed"),
    ("Other", "~4.0", "~", "—", "Diversified"),
]
tbl2 = s.shapes.add_table(len(seg_data),5,Inches(0.5),Inches(2.0),Inches(12.3),Inches(2.8)).table
for ci in range(5): tbl2.columns[ci].width = Inches(12.3/5)
for ri,row in enumerate(seg_data):
    for ci,val in enumerate(row):
        c = tbl2.cell(ri,ci); c.text = val
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(10); p.alignment = PP_ALIGN.CENTER
            if ri==0: p.font.bold=True; p.font.color.rgb=BLUE
            else: p.font.color.rgb=WHITE
        c.fill.solid(); c.fill.fore_color.rgb = DARK_CARD if ri%2==0 else RGBColor(0x16,0x16,0x2A)

add_multi(s, Inches(0.5), Inches(5.2), Inches(12), Inches(2), [
    "STRATEGIC PIVOT: From Apple-dependent manufacturer to diversified precision tech platform",
    "  Leoni: 1st-year profitability improvement, global auto customer base + production network",
    "  AI Infra: High-speed interconnect, optical, liquid cooling, power management for AI servers",
    "  OpenAI Partnership: Future device production — reduces Apple concentration",
    "  HKD IPO: Seeking ~$3B in Hong Kong — global capital market expansion",
], 12, GRAY, True)

# SLIDE 6: Valuation & Ratios
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), BLUE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Valuation & Key Ratios", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A,0x2A,0x4A))

r_data = [
    ("Metric", "FY2025", "Metric", "FY2025"),
    ("P/E (Trailing)", "17.1x", "P/E (Forward)", "~14.5x"),
    ("P/S Ratio", "0.86x", "P/B Ratio", "3.3x"),
    ("PEG Ratio", "0.67", "EV/EBITDA", "14.4x"),
    ("ROE", "19.1%", "ROA", "3.3%"),
    ("Net Margin", "5.00%", "Operating Margin", "4.11%"),
    ("Debt/Equity", "101.5%", "Current Ratio", "1.25"),
    ("Cash", "RMB 60.6B", "Book Value / Share", "RMB 10.02"),
    ("Div Yield", "0.9%", "Insider Ownership", "38.8%"),
]
tbl3 = s.shapes.add_table(len(r_data),4,Inches(0.5),Inches(2.0),Inches(12.3),Inches(4.5)).table
for ci in range(4): tbl3.columns[ci].width = Inches(12.3/4)
for ri,row in enumerate(r_data):
    for ci,val in enumerate(row):
        c = tbl3.cell(ri,ci); c.text = val
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(10); p.alignment = PP_ALIGN.CENTER
            if ri==0: p.font.bold=True; p.font.color.rgb=BLUE
            else: p.font.color.rgb=WHITE
        c.fill.solid(); c.fill.fore_color.rgb = DARK_CARD if ri%2==0 else RGBColor(0x16,0x16,0x2A)

# SLIDE 7: SWOT / Verdict
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), BLUE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Analysis Summary & Outlook", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A,0x2A,0x4A))

add_multi(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5), [
    "STRENGTHS",
    "Record revenue +23.6%, profit +24.2%",
    "Auto: Leoni makes Luxshare a global Tier-1",
    "AI infra: Data center +33.8%, margins improving",
    "Deep Apple + OpenAI partnerships",
    "Strong ROE 19.1%, cash RMB 60.6B",
    "GS target RMB 106 vs current ~60",
    "PEG 0.67 — undervalued for growth rate",
], 12, GREEN, True)

add_multi(s, Inches(7), Inches(2.0), Inches(5.8), Inches(5), [
    "RISKS & CHALLENGES",
    "Apple concentration (~80% consumer rev.)",
    "Debt/Equity 101.5% post-Leoni",
    "Negative FCF (-RMB 579M)",
    "Low net margin (~5%)",
    "US-China trade / tariff uncertainty",
    "Share price -16% over 52 weeks",
    "CapEx heavy (~12% of revenue)",
], 12, RED, True)

add_shape(s, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.03), RGBColor(0x2A,0x2A,0x4A))
add_text(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.8),
    "OUTLOOK: Luxshare is transforming from an Apple-centric assembler into a diversified precision tech platform spanning auto, AI infrastructure, and consumer electronics. The HKD IPO and OpenAI partnership are key catalysts. At 14.5x forward P/E with 24% earnings growth, valuation remains compelling.",
    14, BLUE, True, PP_ALIGN.CENTER)

# SLIDE 8: Thank You
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
add_shape(s, Inches(5.5), Inches(2.8), Inches(2.333), Inches(0.06), BLUE)
add_text(s, Inches(1), Inches(3.2), Inches(11), Inches(1), "Thank You", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.2), Inches(11), Inches(0.5), "Luxshare ICT Financial Analysis  FY2025", 20, GRAY, False, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.2), Inches(11), Inches(0.5), "Data sourced from Luxshare ICT 2025 Annual Report & SZSE filings", 13, SUBDUED, False, PP_ALIGN.CENTER)

prs.save(r"C:\Users\xiaomi\OneDrive\Desktop\Github\Luxshare_Analysis_FY2025.pptx")
print("Created: Luxshare_Analysis_FY2025.pptx")
