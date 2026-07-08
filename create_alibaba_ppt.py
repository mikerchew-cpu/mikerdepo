from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x0F, 0x0F, 0x1A)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
GREEN = RGBColor(0x4E, 0xCC, 0xA3)
PURPLE = RGBColor(0xB3, 0x88, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xB8)
DARK_CARD = RGBColor(0x1A, 0x1A, 0x2E)
SUBDUED = RGBColor(0x7A, 0x7A, 0x9A)

def add_bg(slide, color=DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, size=14, color=GRAY, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(4)
    return txBox

# SLIDE 1: Cover
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), ORANGE)
add_text(s, Inches(1), Inches(1.5), Inches(11), Inches(1), "ALIBABA GROUP", 60, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.4), Inches(11), Inches(0.6), "Financial Analysis   FY2026 (ended March 31, 2026)", 32, ORANGE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.4), Inches(11), Inches(0.5), "NYSE: BABA  |  HKEX: 9988  |  Global Technology Leader", 16, GRAY, False, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.4), "E-Commerce  |  Cloud Computing  |  AI   |   124,000+ employees", 13, SUBDUED, False, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.4), "Market Cap ~$245B  |  Share Price ~$96 (Jul 2026)", 14, PURPLE, False, PP_ALIGN.CENTER)

# SLIDE 2: Executive Summary
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), ORANGE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Executive Summary", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A, 0x2A, 0x4A))
add_multiline(s, Inches(0.5), Inches(2.0), Inches(12), Inches(4.5), [
    "Alibaba Group reported FY2026 revenue of $148.4B, with Cloud Intelligence revenue accelerating to +40% growth in Q4 driven by AI commercialization at scale. Net income was $14.81B, down from $17.36B in FY2025 due to heavy AI infrastructure investment.",
    "",
    "KEY HIGHLIGHTS:",
    "  Revenue: $148.4B (+7.5% YoY)               |   Cloud Growth: +40% (Q4 FY26)",
    "  Net Income: $14.81B (-14.7% YoY)            |   AI Revenue: Triple-digit growth 11th quarter",
    "  Share Buybacks: $11B+ (FY25)                |   Dividend: $4.6B returned",
    "  Cash Position: $19.07B                       |   Employees: 124,320",
    "",
    "CEO Eddie Wu: \"Alibaba's full-stack AI investments have progressed from incubation to commercialization at scale.\"",
], 13, GRAY)

# SLIDE 3: Income Statement
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), ORANGE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Income Statement Highlights", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A, 0x2A, 0x4A))

tbl_data = [
    ("Metric", "FY2026", "FY2025", "Change", "Trend"),
    ("Total Revenue", "$148.40B", "$138.0B*", "+7.5%", "Up"),
    ("Gross Profit", "$59.08B", "$55.6B*", "+6.3%", "Up"),
    ("Operating Income", "$7.27B", "$10.3B*", "-29.4%", "Down"),
    ("Net Income", "$14.81B", "$17.36B", "-14.7%", "Down"),
    ("Diluted EPS", "$6.02", "$7.05", "-14.6%", "Down"),
    ("Revenue (CNY)", "~1.07T", "996.35B", "+7.5%*", "Up"),
]
rows, cols = len(tbl_data), 5
tbl = s.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(12.3), Inches(3.5)).table
for ci in range(cols):
    tbl.columns[ci].width = Inches(12.3 / cols)
for ri, row in enumerate(tbl_data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = ORANGE
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_CARD if ri % 2 == 0 else RGBColor(0x16, 0x16, 0x2A)

add_text(s, Inches(0.5), Inches(5.8), Inches(8), Inches(0.4), "* FY2025 revenue estimated; FY2026 CNY revenue estimated at ~1.07T at ~7.5% growth", 9, SUBDUED, False)

# SLIDE 4: KPI Cards
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), ORANGE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Key Performance Indicators", 28, WHITE, True)

kpis = [
    ("Revenue FY2026", "$148.4B", "+7.5% YoY", GREEN),
    ("Cloud Growth (Q4)", "+40%", "AI-driven acceleration", GREEN),
    ("Net Income", "$14.81B", "-14.7% YoY", ORANGE),
    ("Buybacks (FY25)", "$11B+", "Aggressive returns", GREEN),
    ("Cash Position", "$19.07B", "Strong liquidity", GREEN),
    ("AI Revenue", "Triple-digit", "11 straight quarters", PURPLE),
    ("Dividends (FY25)", "$4.6B", "$1.00/ADS", GREEN),
    ("Employees", "124,320", "Global workforce", PURPLE),
]
for i, (label, value, change, clr) in enumerate(kpis):
    col = i % 4; row = i // 4
    x = Inches(0.5 + col * 3.15); y = Inches(2.0 + row * 2.2)
    card = add_shape(s, x, y, Inches(2.9), Inches(1.3), DARK_CARD)
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(2.5), Inches(0.3), label, 10, SUBDUED, False)
    add_text(s, x + Inches(0.2), y + Inches(0.35), Inches(2.5), Inches(0.5), value, 22, WHITE, True)
    add_text(s, x + Inches(0.2), y + Inches(0.85), Inches(2.5), Inches(0.3), change, 11, clr, False)

# SLIDE 5: Segment Analysis
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), ORANGE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Segment & Strategic Analysis", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A, 0x2A, 0x4A))

seg_data = [
    ("Segment", "Key Highlight", "Growth", "Strategic Importance"),
    ("Cloud Intelligence", "External revenue acceleration", "+40% (Q4)", "AI commercialization at scale"),
    ("AI Products", "30% of Cloud revenue", "Triple-digit x11", "Core growth engine"),
    ("Taobao & Tmall", "Domestic e-commerce growth", "Accelerating", "Cash cow, AI integration"),
    ("AIDC", "International e-commerce", "Strong growth", "Global expansion"),
    ("Cainiao / Local Svcs", "Restructured reporting", "Narrowing loss", "Ecosystem support"),
]
rows2, cols2 = len(seg_data), 4
tbl2 = s.shapes.add_table(rows2, cols2, Inches(0.5), Inches(2.0), Inches(12.3), Inches(3.0)).table
for ci in range(cols2):
    tbl2.columns[ci].width = Inches(12.3 / cols2)
for ri, row in enumerate(seg_data):
    for ci, val in enumerate(row):
        cell = tbl2.cell(ri, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = ORANGE
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_CARD if ri % 2 == 0 else RGBColor(0x16, 0x16, 0x2A)

add_multiline(s, Inches(0.5), Inches(5.3), Inches(12), Inches(2), [
    "STRATEGIC PIVOT: \"User First, AI Driven\"",
    "  Qwen LLM: Leadership in reasoning & coding; multimodal (video, world models); 300M+ open-source downloads",
    "  AI Agents: Enterprise AI agents for office & coding; consumer AI integrated into Qwen app with e-commerce",
    "  Non-core Divestiture: Sold SunArt, Intime to sharpen focus on core e-commerce + cloud",
    "  T-Head (AI Chip): IPO reportedly planned — potential high valuation catalyst",
], 12, GRAY, True)

# SLIDE 6: Balance Sheet & Metrics
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), ORANGE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Balance Sheet & Key Metrics", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A, 0x2A, 0x4A))

bs_data = [
    ("Item", "FY2026", "Item", "FY2026"),
    ("Total Assets", "$244.4B", "Cash & Equivalents", "$19.07B"),
    ("Total Liabilities", "$113.55B", "Long-term Investments", "$58.7B"),
    ("Shareholders' Equity", "~$130.9B", "Book Value / ADS", "$8.32"),
    ("Gross Margin", "~39.8%", "Operating Margin", "~4.9%"),
    ("Net Margin", "~10.0%", "P/E Ratio", "~16.5x"),
    ("P/S Ratio", "~1.65x", "Free Cash Flow (Q4)", "$3.74B"),
]
rows3, cols3 = len(bs_data), 4
tbl3 = s.shapes.add_table(rows3, cols3, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.0)).table
for ci in range(cols3):
    tbl3.columns[ci].width = Inches(12.3 / cols3)
for ri, row in enumerate(bs_data):
    for ci, val in enumerate(row):
        cell = tbl3.cell(ri, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = ORANGE
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_CARD if ri % 2 == 0 else RGBColor(0x16, 0x16, 0x2A)

# SLIDE 7: SWOT / Verdict
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, Inches(0), Inches(1.1), Inches(0.06), Inches(0.5), ORANGE)
add_text(s, Inches(0.3), Inches(1.1), Inches(12), Inches(0.5), "Analysis Summary & Outlook", 28, WHITE, True)
add_shape(s, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.015), RGBColor(0x2A, 0x2A, 0x4A))

add_multiline(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5), [
    "STRENGTHS",
    "Cloud +40% growth — AI commercialization at scale",
    "Qwen LLM leadership (reasoning, coding, multimodal)",
    "Strong balance sheet: $19B cash",
    "$15.6B+ returned to shareholders (FY25)",
    "E-commerce ecosystem: TTG + AIDC global reach",
    "AI chip T-Head spinoff — potential catalyst",
    "300M+ open-source model downloads globally",
], 12, GREEN, True)

add_multiline(s, Inches(7), Inches(2.0), Inches(5.8), Inches(5), [
    "RISKS & CHALLENGES",
    "Net income -14.7% — heavy AI investment costs",
    "Operating margin compressed to ~4.9%",
    "Regulatory overhang in China",
    "Competition: PDD, JD.com, ByteDance, Tencent",
    "Asset base shrinking ($244B vs $255B prior)",
    "China macro headwinds / consumer slowdown",
    "Liabilities up 15.4%, debt/equity rising",
], 12, ORANGE, True)

add_shape(s, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.03), RGBColor(0x2A, 0x2A, 0x4A))
add_text(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.8),
    "OUTLOOK: Alibaba is pivoting from e-commerce giant to AI-first technology leader. AI monetization, T-Head IPO, and international e-commerce profitability are key catalysts for FY2027.",
    14, PURPLE, True, PP_ALIGN.CENTER)

# SLIDE 8: Thank You
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, Inches(5.5), Inches(2.8), Inches(2.333), Inches(0.06), ORANGE)
add_text(s, Inches(1), Inches(3.2), Inches(11), Inches(1), "Thank You", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.2), Inches(11), Inches(0.5), "Alibaba Group Financial Analysis  FY2026", 20, GRAY, False, PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.2), Inches(11), Inches(0.5), "Data sourced from Alibaba Group earnings release (May 2026) & SEC Form 20-F", 13, SUBDUED, False, PP_ALIGN.CENTER)

prs.save(r"C:\Users\xiaomi\OneDrive\Desktop\Github\Alibaba_Analysis_FY2026.pptx")
print("Created: Alibaba_Analysis_FY2026.pptx")
